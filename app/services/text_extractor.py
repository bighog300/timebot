import logging
from pathlib import Path
from typing import Optional, Tuple

logger = logging.getLogger(__name__)


class TextExtractor:
    def extract(self, file_path: Path, file_type: str) -> Tuple[Optional[str], Optional[int], Optional[int]]:
        """Return (text, page_count, word_count)."""
        try:
            ft = file_type.lower().lstrip(".")
            if ft == "pdf":
                return self._extract_pdf(file_path)
            elif ft in ("docx", "doc"):
                return self._extract_docx(file_path)
            elif ft in ("xlsx", "xls"):
                return self._extract_xlsx(file_path)
            elif ft in ("pptx", "ppt"):
                return self._extract_pptx(file_path)
            elif ft == "txt":
                return self._extract_txt(file_path)
            elif ft in ("jpg", "jpeg", "png", "gif", "tiff", "bmp"):
                return self._extract_image(file_path)
            else:
                logger.warning("Unsupported file type: %s", ft)
                return None, None, None
        except Exception as e:
            logger.error("Text extraction failed for %s: %s", file_path, e)
            return None, None, None

    def _extract_pdf(self, file_path: Path) -> Tuple[Optional[str], int, Optional[int]]:
        try:
            import PyPDF2

            parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                page_count = len(reader.pages)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        parts.append(t)

            text = "\n".join(parts)
            if not text.strip() and page_count > 0:
                text = self._pdf_ocr(file_path)

            return text, page_count, len(text.split()) if text else 0
        except Exception as e:
            logger.error("PDF extraction error: %s", e)
            return None, 0, None

    def _pdf_ocr(self, file_path: Path) -> str:
        try:
            from pdf2image import convert_from_path
            import pytesseract

            images = convert_from_path(str(file_path), dpi=200, first_page=1, last_page=5)
            return "\n".join(pytesseract.image_to_string(img) for img in images)
        except Exception as e:
            logger.warning("PDF OCR fallback failed: %s", e)
            return ""

    def _extract_docx(self, file_path: Path) -> Tuple[Optional[str], None, Optional[int]]:
        try:
            from docx import Document

            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text, None, len(text.split())
        except Exception as e:
            logger.error("DOCX extraction error: %s", e)
            return None, None, None

    def _extract_xlsx(self, file_path: Path) -> Tuple[Optional[str], Optional[int], Optional[int]]:
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        rows.append(line)
            text = "\n".join(rows)
            return text, len(wb.worksheets), len(text.split())
        except Exception as e:
            logger.error("XLSX extraction error: %s", e)
            return None, None, None

    def _extract_pptx(self, file_path: Path) -> Tuple[Optional[str], Optional[int], Optional[int]]:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            text = "\n".join(parts)
            return text, len(prs.slides), len(text.split())
        except Exception as e:
            logger.error("PPTX extraction error: %s", e)
            return None, None, None

    def _extract_txt(self, file_path: Path) -> Tuple[Optional[str], int, Optional[int]]:
        try:
            text = file_path.read_text(encoding="utf-8", errors="replace")
            return text, 1, len(text.split())
        except Exception as e:
            logger.error("TXT extraction error: %s", e)
            return None, None, None

    def _extract_image(self, file_path: Path) -> Tuple[Optional[str], int, Optional[int]]:
        try:
            import pytesseract
            from PIL import Image

            text = pytesseract.image_to_string(Image.open(file_path))
            return text, 1, len(text.split())
        except Exception as e:
            logger.error("Image OCR error: %s", e)
            return None, None, None


text_extractor = TextExtractor()
